"""FR-NRM-05/06 (PRD §6.2): OCR, EXIF, and document-text-extraction
primitives for item 6's media pipeline (app/capture/media_pipeline.py). Kept
separate from that orchestration module the same way app/llm/client.py
(the model call) is kept separate from app/ledger/extractor.py (what calls
it and what it's used for) — this module only knows how to turn bytes into
text/metadata, never touches the database or StorageBackend.
"""

import asyncio
import io
import logging
import subprocess
from typing import Protocol

logger = logging.getLogger("cue.capture.media")


class OCRClient(Protocol):
    async def extract_text(self, image_bytes: bytes) -> str: ...


class OCRUnavailableError(Exception):
    """Raised at construction, not first use — a missing OCR dependency
    should fail loudly and immediately, the same posture
    app/capture/adapters/errors.py's CaptureConfigError already establishes
    for a missing credential."""


class TesseractOCRClient:
    """The real, working OCR path in *this* environment today — Tesseract
    (a system binary already present here; `chi_sim`/`chi_tra` language
    data required for Chinese) via pytesseract, a thin, dependency-light
    wrapper (no heavy ML framework of its own). CUE-Tech-Stack.md §2.4 names
    PaddleOCR as the actual *production* choice ("substantially better than
    Tesseract on Chinese... vendor documents are photographed, not
    scanned") — this class is not a claim Tesseract is equivalent, it is the
    same "genuine, real FOSS substitute while the named production choice's
    (much heavier) dependency isn't installed" posture CUE-PRD.md §9.3a
    already established for Mattermost/IMAP-SMTP/Nextcloud, applied here to
    a library instead of a channel integration. See PaddleOCRClient below
    for the interface CUE should actually run once paddlepaddle is worth
    installing.
    """

    def __init__(self, languages: str = "eng+chi_sim"):
        try:
            import pytesseract
        except ImportError as e:
            raise OCRUnavailableError("pytesseract is not installed") from e
        self._pytesseract = pytesseract
        self._languages = languages

    def _extract_sync(self, image_bytes: bytes) -> str:
        from PIL import Image

        image = Image.open(io.BytesIO(image_bytes))
        return self._pytesseract.image_to_string(image, lang=self._languages)

    async def extract_text(self, image_bytes: bytes) -> str:
        return await asyncio.to_thread(self._extract_sync, image_bytes)


class PaddleOCRClient:
    """CUE-Tech-Stack.md §2.4's own named production choice. Not installed
    in this environment — paddleocr pulls in the full paddlepaddle
    deep-learning framework plus a model-weight download on first use, a
    materially heavier install than this session's time/environment budget
    covers. Lazy-imported at construction, not at module import — same
    "fail loudly and immediately, not partway into a live call" posture
    CaptureConfigError already establishes — so importing this *module*
    never requires paddleocr at all, only actually constructing this class
    does. Code-complete, dependency-blocked: the same honest posture this
    codebase already gives WhatsApp/WeChat/Graph for a missing credential,
    applied here to a missing (heavy, optional) library instead.
    """

    def __init__(self, languages: str = "ch"):
        try:
            from paddleocr import PaddleOCR
        except ImportError as e:
            raise OCRUnavailableError(
                "paddleocr is not installed in this environment — see this class's own "
                "docstring; app/capture/media.py's TesseractOCRClient is the real substitute"
            ) from e
        self._engine = PaddleOCR(lang=languages)

    def _extract_sync(self, image_bytes: bytes) -> str:
        import numpy as np
        from PIL import Image

        image = np.array(Image.open(io.BytesIO(image_bytes)).convert("RGB"))
        result = self._engine.ocr(image)
        lines = [line[1][0] for block in (result or []) for line in block]
        return "\n".join(lines)

    async def extract_text(self, image_bytes: bytes) -> str:
        return await asyncio.to_thread(self._extract_sync, image_bytes)


def get_default_ocr_client() -> OCRClient:
    """Prefers the real production choice (PaddleOCR) if its dependency
    happens to be installed, falls back to the real Tesseract substitute
    otherwise — checked fresh on every call (construction is cheap for
    Tesseract; PaddleOCR's own model loading cost is paid once inside
    whichever caller holds onto the returned client, same as any other
    provider-pluggable factory in this codebase, e.g. app/llm/factory.py's
    get_client)."""
    try:
        return PaddleOCRClient()
    except OCRUnavailableError:
        return TesseractOCRClient()


def extract_exif(image_bytes: bytes) -> dict | None:
    """FR-NRM-06: timestamp and geolocation from an onsite photograph's own
    EXIF data — real, via Pillow (already a transitive dependency in this
    codebase, no new install needed). Returns None for an image with no
    EXIF block at all (a screenshot, a re-saved/stripped image, a non-image
    or corrupt payload) — the common, unremarkable case for anything that
    isn't straight off a phone camera, not an error condition."""
    from PIL import ExifTags, Image

    try:
        image = Image.open(io.BytesIO(image_bytes))
        raw_exif = image.getexif()
    except Exception:
        logger.info("EXIF extraction skipped — not a decodable image", exc_info=True)
        return None
    if not raw_exif:
        return None

    tags = {ExifTags.TAGS.get(tag_id, tag_id): value for tag_id, value in raw_exif.items()}
    result: dict = {}
    if "DateTimeOriginal" in tags:
        result["datetime_original"] = str(tags["DateTimeOriginal"])
    elif "DateTime" in tags:
        result["datetime_original"] = str(tags["DateTime"])

    gps_ifd = raw_exif.get_ifd(ExifTags.IFD.GPSInfo)
    if gps_ifd:
        gps_tags = {ExifTags.GPSTAGS.get(tag_id, tag_id): value for tag_id, value in gps_ifd.items()}
        lat = _gps_to_decimal(gps_tags.get("GPSLatitude"), gps_tags.get("GPSLatitudeRef"))
        lon = _gps_to_decimal(gps_tags.get("GPSLongitude"), gps_tags.get("GPSLongitudeRef"))
        if lat is not None and lon is not None:
            result["latitude"] = lat
            result["longitude"] = lon

    return result or None


def _gps_to_decimal(dms, ref) -> float | None:
    if dms is None or ref is None:
        return None
    degrees, minutes, seconds = (float(v) for v in dms)
    decimal = degrees + minutes / 60 + seconds / 3600
    return -decimal if ref in ("S", "W") else decimal


def extract_pdf_text(pdf_bytes: bytes) -> str:
    """FR-NRM-05's PDF half — poppler's `pdftotext` (already a system
    binary in this environment), stdin-to-stdout (`-` for both), no temp
    file needed. CUE-Tech-Stack.md §2.4 names Docling as the real
    production choice for layout/table fidelity (relevant for spec-claim
    extraction specifically) — pdftotext extracts plain reading-order text
    only, no table structure, the same real-but-lighter-substitute posture
    TesseractOCRClient's own docstring explains, applied here to PDF/Office
    parsing instead of image OCR."""
    result = subprocess.run(
        ["pdftotext", "-", "-"], input=pdf_bytes, capture_output=True, check=True
    )
    return result.stdout.decode("utf-8", errors="replace")


def extract_docx_text(docx_bytes: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(docx_bytes))
    return "\n".join(p.text for p in document.paragraphs if p.text)


def extract_pptx_text(pptx_bytes: bytes) -> str:
    """Reuses python-pptx, already a dependency (PPTX report generation,
    app/reports/) — no new install for this half of FR-NRM-05."""
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(pptx_bytes))
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def extract_xlsx_text(xlsx_bytes: bytes) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True, read_only=True)
    lines = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines)


# Maps a filename's extension to the extractor above — the dispatch key
# item 6's pipeline (and FR-DOC-09's drift-check reuse, item 12) uses rather
# than sniffing content-type, since a filename is what every adapter's
# RawCapturedMedia already reliably carries.
_DOCUMENT_EXTRACTORS = {
    ".pdf": extract_pdf_text,
    ".docx": extract_docx_text,
    ".pptx": extract_pptx_text,
    ".xlsx": extract_xlsx_text,
}


def extract_document_text(filename: str | None, data: bytes) -> str | None:
    """FR-NRM-05's PDF/Office half. Returns None for an unrecognised
    extension (e.g. a plain .txt, or a filename-less attachment) — not
    every document kind needs its own bespoke extractor to be a correct,
    honest implementation of "PDF and Office documents", the two kinds
    FR-NRM-05's own wording names."""
    if not filename or "." not in filename:
        return None
    ext = filename[filename.rindex("."):].lower()
    extractor = _DOCUMENT_EXTRACTORS.get(ext)
    if extractor is None:
        return None
    try:
        return extractor(data)
    except Exception:
        logger.warning("document text extraction failed for %r", filename, exc_info=True)
        return None
