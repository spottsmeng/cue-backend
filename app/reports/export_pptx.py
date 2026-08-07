"""FR-RPT-07: PowerPoint export, against a placeholder-branded template
(Prompt 8's own instruction — real Pico branding isn't available; the swap
point is app/reports/templates.py's registry, read via `model.template`
below, not anything hardcoded in this module). Reads the same
ReportRenderModel app/reports/export_html.py does — WHAT TO BUILD #3's "one
composition path feeding both export formats."
"""

import io
import logging

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.reports.render import ReportRenderModel, RenderSection, RenderTable

logger = logging.getLogger("cue.reports.export_pptx")

_SLIDE_BLANK = 6  # python-pptx's blank layout index in the default template


def _color(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


async def _fetch_image(url: str) -> bytes | None:
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            response = await client.get(url)
            response.raise_for_status()
            return response.content
    except httpx.HTTPError as e:  # a signed URL can legitimately 404/expire between compose and export
        logger.warning("failed to fetch report image %s: %s", url, e)
        return None


def _add_title_slide(prs: Presentation, model: ReportRenderModel) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_SLIDE_BLANK])
    primary = _color(model.template["primary_color"])

    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.4))
    logo_run = logo_box.text_frame.paragraphs[0].add_run()
    logo_run.text = model.template["logo_text"]
    logo_run.font.size = Pt(11)
    logo_run.font.italic = True
    logo_run.font.color.rgb = _color(model.template["accent_color"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = model.project_name
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    title_run.font.color.rgb = primary

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.6))
    subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
    subtitle_run.text = f"Living WIP report — generated {model.generated_at}"
    subtitle_run.font.size = Pt(14)


def _add_heading(slide, model: ReportRenderModel, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.6))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = _color(model.template["primary_color"])


def _add_facts(slide, facts: list[tuple[str, str]], top: float) -> float:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(top), Inches(9.2), Inches(0.35 * len(facts) + 0.2))
    tf = box.text_frame
    for i, (label, value) in enumerate(facts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{label}: {value}"
        run.font.size = Pt(14)
    return top + 0.35 * len(facts) + 0.4


def _add_table(prs: Presentation, slide, model: ReportRenderModel, table: RenderTable, top: float) -> None:
    label_box = slide.shapes.add_textbox(Inches(0.4), Inches(top), Inches(9.2), Inches(0.3))
    run = label_box.text_frame.paragraphs[0].add_run()
    run.text = table.title
    run.font.size = Pt(14)
    run.font.bold = True

    rows = table.rows or [["—"] * len(table.headers)]
    n_rows, n_cols = len(rows) + 1, len(table.headers)
    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.4), Inches(top + 0.35), Inches(9.2), Inches(0.3 * n_rows)
    )
    gfx_table = shape.table
    for c, header in enumerate(table.headers):
        cell = gfx_table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = gfx_table.cell(r, c)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)


async def _add_section_slides(prs: Presentation, model: ReportRenderModel, section: RenderSection) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[_SLIDE_BLANK])
    _add_heading(slide, model, section.heading)
    top = 1.0
    if section.facts:
        top = _add_facts(slide, section.facts, top)
    if section.note:
        note_box = slide.shapes.add_textbox(Inches(0.4), Inches(top), Inches(9.2), Inches(0.3))
        run = note_box.text_frame.paragraphs[0].add_run()
        run.text = section.note
        run.font.size = Pt(11)
        run.font.italic = True
        top += 0.4

    # Images (project overview's visual references, FR-RPT-03) get their own
    # slide each — a table/text slide has no room for print-resolution
    # imagery alongside it.
    for caption, url in section.images:
        image_bytes = await _fetch_image(url)
        if image_bytes is None:
            continue
        image_slide = prs.slides.add_slide(prs.slide_layouts[_SLIDE_BLANK])
        _add_heading(image_slide, model, f"{section.heading}: {caption}")
        image_slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(1.0), Inches(1.2), height=Inches(5.0))

    for table in section.tables:
        if top > 6.5:  # ran out of room on this slide — start a fresh one for the remaining tables
            slide = prs.slides.add_slide(prs.slide_layouts[_SLIDE_BLANK])
            _add_heading(slide, model, f"{section.heading} (cont'd)")
            top = 1.0
        _add_table(prs, slide, model, table, top)
        top += 0.35 * (len(table.rows) + 1) + 0.6


async def render_pptx(model: ReportRenderModel) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, model)
    for section in model.sections:
        await _add_section_slides(prs, model, section)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
