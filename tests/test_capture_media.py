"""FR-NRM-05/06: app/capture/media.py's OCR, EXIF and document-text-
extraction primitives — real Tesseract, real Pillow EXIF, real python-docx/
python-pptx/openpyxl round trips, real pdftotext against a real (if
hand-minimal) PDF. Not mocked: this module's whole job is turning real bytes
into real text, so a mock would prove nothing about whether it works.
"""

import io

import docx
import openpyxl
import pytest
from PIL import Image
from pptx import Presentation

from app.capture.media import (
    OCRUnavailableError,
    TesseractOCRClient,
    extract_docx_text,
    extract_document_text,
    extract_exif,
    extract_pdf_text,
    extract_pptx_text,
    extract_xlsx_text,
    get_default_ocr_client,
)

_MINIMAL_PDF = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/Resources<</Font<</F1 4 0 R>>>>/MediaBox[0 0 200 200]/Contents 5 0 R>>endobj
4 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
5 0 obj<</Length 58>>
stream
BT /F1 18 Tf 20 100 Td (Hello CUE PDF test) Tj ET
endstream
endobj
xref
0 6
trailer<</Size 6/Root 1 0 R>>
%%EOF"""


def _render_text_image(text: str) -> bytes:
    from PIL import ImageDraw

    image = Image.new("RGB", (600, 100), color="white")
    ImageDraw.Draw(image).text((10, 30), text, fill="black")
    buf = io.BytesIO()
    image.save(buf, format="png")
    return buf.getvalue()


@pytest.mark.asyncio
async def test_tesseract_ocr_reads_rendered_text():
    client = TesseractOCRClient()
    text = await client.extract_text(_render_text_image("DELIVERY CONFIRMED"))
    assert "DELIVERY CONFIRMED" in text


def test_get_default_ocr_client_returns_a_working_client():
    """Whichever engine is actually available in this environment
    (PaddleOCR if installed, Tesseract otherwise) — this must never raise
    OCRUnavailableError itself, since Tesseract is a hard local dependency
    this test suite already requires."""
    client = get_default_ocr_client()
    assert client is not None


def test_extract_exif_returns_none_for_image_with_no_exif():
    assert extract_exif(_render_text_image("no exif here")) is None


def test_extract_exif_returns_none_for_non_image_bytes():
    assert extract_exif(b"not an image at all") is None


def test_extract_exif_reads_datetime_and_gps():
    image = Image.new("RGB", (4, 4), color="red")
    exif = image.getexif()
    exif[0x0132] = "2026:06:22 09:00:00"  # DateTime
    from PIL import ExifTags

    gps_ifd = exif.get_ifd(ExifTags.IFD.GPSInfo)
    gps_ifd[1] = "N"  # GPSLatitudeRef
    gps_ifd[2] = (1.0, 21.0, 0.0)  # GPSLatitude: 1 deg 21 min -> 1.35
    gps_ifd[3] = "E"  # GPSLongitudeRef
    gps_ifd[4] = (103.0, 49.0, 0.0)  # GPSLongitude
    buf = io.BytesIO()
    image.save(buf, format="jpeg", exif=exif)

    result = extract_exif(buf.getvalue())
    assert result is not None
    assert result["datetime_original"] == "2026:06:22 09:00:00"
    assert result["latitude"] == pytest.approx(1.35, abs=1e-6)
    assert result["longitude"] == pytest.approx(103.8166667, abs=1e-4)


def test_extract_pdf_text():
    assert extract_pdf_text(_MINIMAL_PDF).strip() == "Hello CUE PDF test"


def test_extract_docx_text():
    document = docx.Document()
    document.add_paragraph("LED screen install confirmed for Friday.")
    buf = io.BytesIO()
    document.save(buf)
    assert "LED screen install confirmed for Friday." in extract_docx_text(buf.getvalue())


def test_extract_pptx_text():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 1000000, 500000)
    box.text_frame.text = "Graphics test print approval"
    buf = io.BytesIO()
    presentation.save(buf)
    assert "Graphics test print approval" in extract_pptx_text(buf.getvalue())


def test_extract_xlsx_text():
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet["A1"] = "Item"
    sheet["B1"] = "Qty"
    sheet["A2"] = "LED panel"
    sheet["B2"] = 12
    buf = io.BytesIO()
    workbook.save(buf)
    text = extract_xlsx_text(buf.getvalue())
    assert "Item" in text and "LED panel" in text and "12" in text


def test_extract_document_text_dispatches_by_extension():
    document = docx.Document()
    document.add_paragraph("dispatch check")
    buf = io.BytesIO()
    document.save(buf)
    assert "dispatch check" in extract_document_text("quote.docx", buf.getvalue())


def test_extract_document_text_returns_none_for_unknown_extension():
    assert extract_document_text("notes.txt", b"plain text") is None


def test_extract_document_text_returns_none_for_missing_filename():
    assert extract_document_text(None, b"x") is None
